"""Xuất bộ slide giảng (.pptx) từ sổ tay trong docs/sessions/.

Khác với lớp học OpenMAIC (xem scripts/openmaic_export.py), bộ slide này dành cho
NGƯỜI đứng lớp: mỗi slide đều có ghi chú thuyết minh trong phần speaker notes, và
mỗi câu trắc nghiệm có một slide đáp án riêng để bấm sau khi lớp đã chọn.

Nguồn nội dung — không viết lại chữ nào bằng LLM:

    docs/sessions/session-*.md      nội dung giảng
    docs/95-ngan-hang-cau-hoi.md    trắc nghiệm và chủ đề thảo luận

Dùng (python-pptx nạp tạm, không thêm phụ thuộc vào dự án):

    uv run --with python-pptx python scripts/slide_giang.py
    uv run --with python-pptx python scripts/slide_giang.py --bai 1 3

File .pptx nằm ở reports/slides/.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from openmaic_export import (  # noqa: E402
    MUC_BAI_TAP,
    MUC_BO_QUA,
    MUC_KIEM_TRA,
    MUC_MUC_LUC,
    NGUON,
    Bai,
    Muc,
    NoiDungLop,
    TracNghiem,
    _bullet_tu_muc,
    _dan_khoi,
    _loi_giang,
    _ten_file,
    doc_bai,
    doc_ngan_hang,
)

DICH = ROOT / "reports" / "slides"

# Bảng màu lấy từ chính chủ đề khoá học: xanh là test đúng, đỏ là test hỏng.
NEN_TOI = "121A21"
NEN_CODE = "0D141A"
TRANG = "FFFFFF"
CHU = "1B2430"
CHU_PHU = "5B6B7A"
CHU_TOI_PHU = "9DB2BF"
XANH = "2FBF71"
DO = "E5484D"
NEN_LUU_Y = "F1FAF4"
NEN_CANH_BAO = "FDF0F0"
NEN_NHAT = "F4F6F8"

CHU_TIEU_DE = "Cambria"
CHU_THAN = "Calibri"
CHU_MA = "Courier New"

RONG = 13.333
CAO = 7.5
LE = 0.75


def _mau(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _bo_bong(hinh) -> None:
    """Gỡ <p:style> để hình không nhận bóng đổ từ theme — slide phẳng, sạch."""
    hinh.shadow.inherit = False
    kieu = hinh._element.find(qn("p:style"))
    if kieu is not None:
        hinh._element.remove(kieu)


def _nen(slide, mau: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _mau(mau)


def _tren_can_giua(cao_khoi: float, tren_toi_thieu: float = 1.9) -> float:
    """Khối ngắn thì đẩy xuống giữa vùng nội dung thay vì để trống nửa dưới."""
    vung_tren, vung_duoi = tren_toi_thieu, CAO - 0.7
    return max(tren_toi_thieu, vung_tren + (vung_duoi - vung_tren - cao_khoi) / 2)


def _hop_chu(
    slide,
    text,
    *,
    trai,
    tren,
    rong,
    cao,
    co_chu=16,
    mau=CHU,
    dam=False,
    font=CHU_THAN,
    can=PP_ALIGN.LEFT,
    gian_dong=1.25,
    cach_doan=6,
    neo=MSO_ANCHOR.TOP,
):
    """text: chuỗi, hoặc list dòng, hoặc list (chuỗi, dict ghi đè kiểu chữ)."""
    hop = slide.shapes.add_textbox(Inches(trai), Inches(tren), Inches(rong), Inches(cao))
    khung = hop.text_frame
    khung.word_wrap = True
    khung.vertical_anchor = neo
    khung.margin_left = khung.margin_right = 0
    khung.margin_top = khung.margin_bottom = 0

    dong = text if isinstance(text, list) else [text]
    for i, muc in enumerate(dong):
        noi_dung, ghi_de = muc if isinstance(muc, tuple) else (muc, {})
        doan = khung.paragraphs[0] if i == 0 else khung.add_paragraph()
        doan.alignment = ghi_de.get("can", can)
        doan.line_spacing = ghi_de.get("gian_dong", gian_dong)
        doan.space_after = Pt(ghi_de.get("cach_doan", cach_doan))
        chay = doan.add_run()
        chay.text = noi_dung
        chay.font.size = Pt(ghi_de.get("co_chu", co_chu))
        chay.font.bold = ghi_de.get("dam", dam)
        chay.font.name = ghi_de.get("font", font)
        chay.font.color.rgb = _mau(ghi_de.get("mau", mau))
    return hop


def _the(slide, *, trai, tren, rong, cao, mau_nen, bo_tron=True, mau_vien=None):
    hinh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if bo_tron else MSO_SHAPE.RECTANGLE,
        Inches(trai),
        Inches(tren),
        Inches(rong),
        Inches(cao),
    )
    hinh.fill.solid()
    hinh.fill.fore_color.rgb = _mau(mau_nen)
    if mau_vien:
        hinh.line.color.rgb = _mau(mau_vien)
        hinh.line.width = Pt(1)
    else:
        hinh.line.fill.background()
    _bo_bong(hinh)
    if bo_tron:
        hinh.adjustments[0] = 0.08
    hinh.text_frame.text = ""
    return hinh


def _vong_so(slide, so: str, *, trai, tren, duong_kinh=0.62, mau_nen=XANH, mau_chu=TRANG):
    """Mô-típ xuyên suốt bộ slide: số thứ tự trong vòng tròn màu."""
    hinh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(trai), Inches(tren), Inches(duong_kinh), Inches(duong_kinh)
    )
    hinh.fill.solid()
    hinh.fill.fore_color.rgb = _mau(mau_nen)
    hinh.line.fill.background()
    _bo_bong(hinh)
    khung = hinh.text_frame
    khung.margin_left = khung.margin_right = khung.margin_top = khung.margin_bottom = 0
    khung.vertical_anchor = MSO_ANCHOR.MIDDLE
    doan = khung.paragraphs[0]
    doan.alignment = PP_ALIGN.CENTER
    chay = doan.add_run()
    chay.text = so
    chay.font.size = Pt(18 if len(so) < 3 else 14)
    chay.font.bold = True
    chay.font.name = CHU_THAN
    chay.font.color.rgb = _mau(mau_chu)
    return hinh


def _ghi_chu(slide, text: str) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _trang(prs) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])


# --------------------------------------------------------------------------- #
# Từng loại slide
# --------------------------------------------------------------------------- #

TEN_KHOA = "Kiểm thử tự động cho ứng dụng AI"


def slide_bia(prs, bai: Bai, mo_dau: str) -> None:
    slide = _trang(prs)
    _nen(slide, NEN_TOI)
    _hop_chu(
        slide,
        f"BÀI {bai.ma}",
        trai=LE,
        tren=1.55,
        rong=6,
        cao=0.45,
        co_chu=18,
        mau=XANH,
        dam=True,
    )
    _hop_chu(
        slide,
        bai.tieu_de,
        trai=LE,
        tren=2.05,
        rong=RONG - 2 * LE - 1.2,
        cao=2.1,
        co_chu=40,
        mau=TRANG,
        dam=True,
        font=CHU_TIEU_DE,
        gian_dong=1.1,
    )
    if bai.muc_tieu:
        _hop_chu(
            slide,
            bai.muc_tieu,
            trai=LE,
            tren=4.35,
            rong=RONG - 2 * LE - 2.2,
            cao=1.0,
            co_chu=19,
            mau=CHU_TOI_PHU,
        )
    _hop_chu(
        slide,
        f"{bai.thoi_luong}   ·   {TEN_KHOA}",
        trai=LE,
        tren=6.35,
        rong=8,
        cao=0.4,
        co_chu=14,
        mau=CHU_TOI_PHU,
    )
    _ghi_chu(slide, mo_dau)


def slide_muc_luc(prs, bai: Bai) -> None:
    muc = next((m for m in bai.muc if m.tieu_de == MUC_MUC_LUC), None)
    bang = next((k for k in muc.khoi if k.loai == "table"), None) if muc else None
    if not bang:
        return
    hang = [h for h in bang.hang[1:] if h and h[0]][:8]

    slide = _trang(prs)
    _nen(slide, TRANG)
    _hop_chu(
        slide,
        "Nội dung buổi học",
        trai=LE,
        tren=0.62,
        rong=9,
        cao=0.8,
        co_chu=34,
        dam=True,
        font=CHU_TIEU_DE,
    )
    tren = 1.75
    cao_hang = min(0.62, (CAO - 2.4) / max(len(hang), 1))
    for i, h in enumerate(hang, start=1):
        _vong_so(slide, str(i), trai=LE, tren=tren + 0.02, duong_kinh=0.42, mau_nen=NEN_NHAT, mau_chu=CHU_PHU)
        _hop_chu(
            slide,
            h[0],
            trai=LE + 0.62,
            tren=tren + 0.05,
            rong=9.0,
            cao=cao_hang,
            co_chu=18,
        )
        if len(h) > 1 and h[1]:
            _hop_chu(
                slide,
                h[1],
                trai=RONG - LE - 1.6,
                tren=tren + 0.07,
                rong=1.6,
                cao=cao_hang,
                co_chu=14,
                mau=CHU_PHU,
                can=PP_ALIGN.RIGHT,
            )
        tren += cao_hang
    _ghi_chu(
        slide,
        "Chặng đường của buổi hôm nay. Bám mốc thời gian để không lệch nhịp — "
        "phần thực hành luôn dài hơn dự kiến.",
    )


def slide_muc(prs, muc: Muc, so: int) -> None:
    slide = _trang(prs)
    _nen(slide, TRANG)
    _vong_so(slide, str(so), trai=LE, tren=0.68)
    _hop_chu(
        slide,
        muc.tieu_de,
        trai=LE + 0.85,
        tren=0.62,
        rong=RONG - 2 * LE - 0.85,
        cao=1.0,
        co_chu=32,
        dam=True,
        font=CHU_TIEU_DE,
        gian_dong=1.05,
    )

    bullet = _bullet_tu_muc(muc)
    luu_y = next((k for k in muc.khoi if k.loai == "callout"), None)
    tren = 1.95
    if bullet:
        cao_the = min(0.55 + 0.62 * len(bullet), 3.9)
        _the(slide, trai=LE, tren=tren, rong=RONG - 2 * LE, cao=cao_the, mau_nen=NEN_NHAT)
        dong = []
        for b, dam in bullet:
            dong.append(
                (
                    b if dam else f"•   {b}",
                    {"dam": dam, "co_chu": 20 if dam else 18, "cach_doan": 12},
                )
            )
        _hop_chu(
            slide,
            dong,
            trai=LE + 0.5,
            tren=tren + 0.32,
            rong=RONG - 2 * LE - 1.0,
            cao=cao_the - 0.5,
            co_chu=18,
        )
        tren += cao_the + 0.35

    if luu_y and tren < CAO - 1.3:
        canh_bao = luu_y.canh_bao
        chu = luu_y.van_ban[:280]
        cao_the = min(0.5 + 0.34 * (len(chu) // 95 + 2), CAO - 0.6 - tren)
        _the(
            slide,
            trai=LE,
            tren=tren,
            rong=RONG - 2 * LE,
            cao=cao_the,
            mau_nen=NEN_CANH_BAO if canh_bao else NEN_LUU_Y,
        )
        _hop_chu(
            slide,
            [
                ("CẢNH BÁO" if canh_bao else "GHI NHỚ", {"co_chu": 12, "dam": True, "mau": DO if canh_bao else XANH}),
                (chu, {"co_chu": 15, "mau": CHU}),
            ],
            trai=LE + 0.45,
            tren=tren + 0.22,
            rong=RONG - 2 * LE - 0.9,
            cao=cao_the - 0.4,
        )

    _ghi_chu(
        slide,
        _loi_giang(muc) or f"Dẫn nhanh qua các ý của mục {muc.tieu_de}, rồi vào phần thực hành.",
    )


def slide_code(prs, tieu_de: str, khoi, muc: Muc, chi_so: int) -> None:
    dong = khoi.van_ban.splitlines()[:18]
    rong_nhat = max((len(d) for d in dong), default=0)
    co_chu = 15 if rong_nhat <= 72 else 13 if rong_nhat <= 92 else 11

    slide = _trang(prs)
    _nen(slide, TRANG)
    _hop_chu(
        slide,
        tieu_de,
        trai=LE,
        tren=0.62,
        rong=RONG - 2 * LE,
        cao=0.9,
        co_chu=30,
        dam=True,
        font=CHU_TIEU_DE,
    )
    cao_the = min(0.7 + len(dong) * (co_chu * 1.55 / 72), CAO - 2.3)
    tren = _tren_can_giua(cao_the)
    _the(slide, trai=LE, tren=tren, rong=RONG - 2 * LE, cao=cao_the, mau_nen=NEN_CODE)
    _hop_chu(
        slide,
        khoi.ngon_ngu.upper(),
        trai=RONG - LE - 1.6,
        tren=tren + 0.17,
        rong=1.4,
        cao=0.3,
        co_chu=11,
        mau=CHU_TOI_PHU,
        can=PP_ALIGN.RIGHT,
    )
    _hop_chu(
        slide,
        [(d if d.strip() else " ", {}) for d in dong],
        trai=LE + 0.42,
        tren=tren + 0.2,
        rong=RONG - 2 * LE - 1.9,
        cao=cao_the - 0.5,
        co_chu=co_chu,
        mau="E6EDF3",
        font=CHU_MA,
        gian_dong=1.25,
        cach_doan=0,
    )
    _ghi_chu(
        slide,
        _dan_khoi(muc, chi_so)
        or _loi_giang(muc)
        or f"Gõ đoạn này cùng lớp, đợi mọi người ra cùng kết quả rồi mới đi tiếp. "
        f"(Mục: {muc.tieu_de})",
    )


def slide_bang(prs, tieu_de: str, khoi, muc: Muc, chi_so: int) -> None:
    hang = [h for h in khoi.hang if h][:9]
    if not hang:
        return
    so_cot = max(len(h) for h in hang)

    slide = _trang(prs)
    _nen(slide, TRANG)
    _hop_chu(
        slide,
        tieu_de,
        trai=LE,
        tren=0.62,
        rong=RONG - 2 * LE,
        cao=0.9,
        co_chu=30,
        dam=True,
        font=CHU_TIEU_DE,
    )
    cao_bang = min(0.5 * len(hang) + 0.2, CAO - 2.4)
    hinh = slide.shapes.add_table(
        len(hang),
        so_cot,
        Inches(LE),
        Inches(_tren_can_giua(cao_bang)),
        Inches(RONG - 2 * LE),
        Inches(cao_bang),
    )
    bang = hinh.table
    bang.first_row = True
    for r, h in enumerate(hang):
        bang.rows[r].height = Emu(int(Inches(cao_bang).emu / len(hang)))
        for c in range(so_cot):
            o = bang.cell(r, c)
            o.text = ""
            o.margin_left = o.margin_right = Inches(0.14)
            o.margin_top = o.margin_bottom = Inches(0.05)
            o.vertical_anchor = MSO_ANCHOR.MIDDLE
            o.fill.solid()
            o.fill.fore_color.rgb = _mau(NEN_TOI if r == 0 else (TRANG if r % 2 else NEN_NHAT))
            doan = o.text_frame.paragraphs[0]
            chay = doan.add_run()
            chay.text = (h[c] if c < len(h) else "")[:110]
            chay.font.size = Pt(14 if r else 13)
            chay.font.bold = r == 0
            chay.font.name = CHU_THAN
            chay.font.color.rgb = _mau(TRANG if r == 0 else CHU)
    _ghi_chu(
        slide,
        _dan_khoi(muc, chi_so)
        or _loi_giang(muc)
        or f"Đọc bảng theo từng hàng, đối chiếu cột đầu với cột cuối. (Mục: {muc.tieu_de})",
    )


def slide_bai_tap(prs, muc: Muc) -> None:
    y = [m for k in muc.khoi if k.loai == "bullets" for m in k.muc][:5]
    if not y:
        return
    slide = _trang(prs)
    _nen(slide, TRANG)
    _hop_chu(
        slide,
        "Bài tập",
        trai=LE,
        tren=0.62,
        rong=8,
        cao=0.9,
        co_chu=32,
        dam=True,
        font=CHU_TIEU_DE,
    )
    _hop_chu(
        slide,
        "Làm trên máy trước — mỗi câu cần một bằng chứng chạy thật.",
        trai=LE,
        tren=1.48,
        rong=10,
        cao=0.4,
        co_chu=15,
        mau=CHU_PHU,
    )
    tren = 2.15
    cao_hang = min(1.0, (CAO - 2.9) / len(y))
    for i, m in enumerate(y, start=1):
        _the(slide, trai=LE, tren=tren, rong=RONG - 2 * LE, cao=cao_hang - 0.12, mau_nen=NEN_NHAT)
        _vong_so(slide, str(i), trai=LE + 0.3, tren=tren + (cao_hang - 0.12 - 0.42) / 2, duong_kinh=0.42)
        _hop_chu(
            slide,
            m[:190],
            trai=LE + 0.95,
            tren=tren + 0.16,
            rong=RONG - 2 * LE - 1.5,
            cao=cao_hang - 0.35,
            co_chu=16,
            neo=MSO_ANCHOR.TOP,
        )
        tren += cao_hang
    _ghi_chu(
        slide,
        "Giao bài tập về nhà. Nhấn mạnh: câu trả lời phải kèm bằng chứng chạy thật, "
        "không phải suy đoán từ việc đọc code.",
    )


def _slide_cau_hoi_khung(prs, tn: TracNghiem, so: int, dap_an: bool):
    slide = _trang(prs)
    _nen(slide, TRANG)
    _vong_so(slide, f"C{so}", trai=LE, tren=0.66, mau_nen=XANH if dap_an else NEN_TOI)
    _hop_chu(
        slide,
        tn.cau_hoi,
        trai=LE + 0.85,
        tren=0.6,
        rong=RONG - 2 * LE - 0.85,
        cao=1.15,
        co_chu=24,
        dam=True,
        font=CHU_TIEU_DE,
        gian_dong=1.15,
    )
    tren = 2.05
    cao_hang = min(0.92, (CAO - (3.3 if dap_an else 2.5)) / max(len(tn.lua_chon), 1))
    for nhan, chu in tn.lua_chon:
        dung = dap_an and nhan in tn.dap_an
        _the(
            slide,
            trai=LE,
            tren=tren,
            rong=RONG - 2 * LE,
            cao=cao_hang - 0.14,
            mau_nen=NEN_LUU_Y if dung else NEN_NHAT,
            mau_vien=XANH if dung else None,
        )
        _hop_chu(
            slide,
            nhan,
            trai=LE + 0.35,
            tren=tren + (cao_hang - 0.14) / 2 - 0.16,
            rong=0.45,
            cao=0.35,
            co_chu=17,
            dam=True,
            mau=XANH if dung else CHU_PHU,
        )
        _hop_chu(
            slide,
            chu,
            trai=LE + 0.95,
            tren=tren + 0.13,
            rong=RONG - 2 * LE - 1.6,
            cao=cao_hang - 0.3,
            co_chu=17,
            dam=dung,
        )
        tren += cao_hang
    return slide, tren


def slide_cau_hoi(prs, tn: TracNghiem, so: int) -> None:
    slide, _ = _slide_cau_hoi_khung(prs, tn, so, dap_an=False)
    _ghi_chu(
        slide,
        f"Cho lớp 60 giây chọn đáp án, hỏi vài người vì sao chọn, rồi mới bấm sang "
        f"slide đáp án. Đáp án đúng: {', '.join(tn.dap_an)}. {tn.giai_thich}",
    )


def slide_dap_an(prs, tn: TracNghiem, so: int) -> None:
    slide, tren = _slide_cau_hoi_khung(prs, tn, so, dap_an=True)
    if tn.giai_thich and tren < CAO - 1.0:
        _hop_chu(
            slide,
            [
                (f"ĐÁP ÁN {', '.join(tn.dap_an)}", {"co_chu": 12, "dam": True, "mau": XANH}),
                (tn.giai_thich, {"co_chu": 15, "mau": CHU_PHU}),
            ],
            trai=LE,
            tren=tren + 0.12,
            rong=RONG - 2 * LE,
            cao=CAO - tren - 0.5,
        )
    _ghi_chu(slide, tn.giai_thich)


def slide_thao_luan(prs, noi_dung: NoiDungLop) -> None:
    if not noi_dung or not noi_dung.thao_luan:
        return
    slide = _trang(prs)
    _nen(slide, NEN_TOI)
    _hop_chu(
        slide,
        "Thảo luận",
        trai=LE,
        tren=0.75,
        rong=8,
        cao=0.9,
        co_chu=34,
        mau=TRANG,
        dam=True,
        font=CHU_TIEU_DE,
    )
    tren = 2.0
    for i, td in enumerate(noi_dung.thao_luan[:2], start=1):
        _vong_so(slide, str(i), trai=LE, tren=tren + 0.05, duong_kinh=0.5)
        _hop_chu(
            slide,
            td.chu_de,
            trai=LE + 0.8,
            tren=tren,
            rong=RONG - 2 * LE - 0.8,
            cao=0.9,
            co_chu=21,
            mau=TRANG,
            dam=True,
            gian_dong=1.15,
        )
        _hop_chu(
            slide,
            td.goi_y,
            trai=LE + 0.8,
            tren=tren + 0.95,
            rong=RONG - 2 * LE - 1.2,
            cao=1.3,
            co_chu=15,
            mau=CHU_TOI_PHU,
        )
        tren += 2.45
    _ghi_chu(
        slide,
        "Chia nhóm 3-4 người, mỗi chủ đề 8 phút, rồi mỗi nhóm nói một câu kết luận. "
        "Không cần đồng thuận — cần lý do.",
    )


def slide_ket(prs, bai: Bai) -> None:
    muc = next((m for m in bai.muc if m.tieu_de == MUC_KIEM_TRA), None)
    khoi = next((k for k in muc.khoi if k.loai == "code"), None) if muc else None
    slide = _trang(prs)
    _nen(slide, NEN_TOI)
    _hop_chu(
        slide,
        "Kiểm tra trước khi sang bài sau",
        trai=LE,
        tren=0.75,
        rong=RONG - 2 * LE,
        cao=0.9,
        co_chu=32,
        mau=TRANG,
        dam=True,
        font=CHU_TIEU_DE,
    )
    if khoi:
        dong = khoi.van_ban.splitlines()[:14]
        co_chu = 15 if max((len(d) for d in dong), default=0) <= 72 else 12
        cao_the = min(0.6 + len(dong) * (co_chu * 1.55 / 72), 4.2)
        _the(slide, trai=LE, tren=2.0, rong=RONG - 2 * LE, cao=cao_the, mau_nen=NEN_CODE)
        _hop_chu(
            slide,
            [(d if d.strip() else " ", {}) for d in dong],
            trai=LE + 0.42,
            tren=2.2,
            rong=RONG - 2 * LE - 1.0,
            cao=cao_the - 0.4,
            co_chu=co_chu,
            mau="E6EDF3",
            font=CHU_MA,
            gian_dong=1.25,
            cach_doan=0,
        )
        _hop_chu(
            slide,
            "Xanh thì sang bài sau. Đỏ thì dừng lại xử lý ngay — đừng mang lỗi sang buổi sau.",
            trai=LE,
            tren=min(2.0 + cao_the + 0.35, CAO - 0.9),
            rong=RONG - 2 * LE,
            cao=0.5,
            co_chu=16,
            mau=CHU_TOI_PHU,
        )
    else:
        _hop_chu(
            slide,
            "Hết bài. Chạy lại bộ test của mình trước khi đóng máy.",
            trai=LE,
            tren=2.2,
            rong=RONG - 2 * LE,
            cao=0.8,
            co_chu=19,
            mau=CHU_TOI_PHU,
        )
    _ghi_chu(
        slide,
        "Chạy đoạn kiểm tra ngay tại lớp, đợi cả lớp báo xanh rồi mới kết thúc buổi.",
    )


def slide_buoc_tiep(prs, bai: Bai) -> None:
    """Slide chốt: chính mục 'Chọn bước tiếp theo' của sổ tay — lối rẽ cho từng người."""
    muc = next((m for m in bai.muc if m.tieu_de == "Chọn bước tiếp theo"), None)
    y = [m for k in (muc.khoi if muc else []) if k.loai == "bullets" for m in k.muc][:4]
    if not y:
        return
    slide = _trang(prs)
    _nen(slide, TRANG)
    _hop_chu(
        slide,
        "Chọn bước tiếp theo",
        trai=LE,
        tren=0.62,
        rong=9,
        cao=0.9,
        co_chu=32,
        dam=True,
        font=CHU_TIEU_DE,
    )
    tren = 1.95
    cao_hang = min(1.15, (CAO - 2.7) / len(y))
    for m in y:
        _the(slide, trai=LE, tren=tren, rong=RONG - 2 * LE, cao=cao_hang - 0.15, mau_nen=NEN_NHAT)
        dau, _, duoi = m.partition("→")
        _hop_chu(
            slide,
            [
                (dau.strip(" —-").strip()[:110], {"co_chu": 17, "dam": True}),
                (duoi.strip()[:110], {"co_chu": 15, "mau": CHU_PHU}) if duoi else ("", {}),
            ],
            trai=LE + 0.5,
            tren=tren + 0.18,
            rong=RONG - 2 * LE - 1.0,
            cao=cao_hang - 0.4,
            cach_doan=2,
        )
        tren += cao_hang
    _ghi_chu(
        slide,
        "Không phải ai cũng đi tiếp cùng một hướng. Đọc to bốn lối rẽ này để mỗi người "
        "tự chọn phần phù hợp với mình.",
    )


# --------------------------------------------------------------------------- #
# Lắp bộ slide
# --------------------------------------------------------------------------- #


def dung_deck(bai: Bai, noi_dung: NoiDungLop | None) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(RONG)
    prs.slide_height = Inches(CAO)

    mo_dau = f"Chào cả lớp. Bài {bai.ma}: {bai.tieu_de}."
    if bai.muc_tieu:
        mo_dau += f" Hết {bai.thoi_luong} này, mục tiêu là: {bai.muc_tieu}"
    if noi_dung and noi_dung.thao_luan:
        mo_dau += f" Mở màn bằng một câu hỏi: {noi_dung.thao_luan[0].chu_de}"

    slide_bia(prs, bai, mo_dau)
    slide_muc_luc(prs, bai)

    so = 0
    for muc in bai.muc:
        if muc.tieu_de in MUC_BO_QUA or muc.tieu_de in {MUC_BAI_TAP, MUC_KIEM_TRA}:
            continue
        so += 1
        slide_muc(prs, muc, so)
        code = [i for i, k in enumerate(muc.khoi) if k.loai == "code"][:3]
        for n, i in enumerate(code, start=1):
            k = muc.khoi[i]
            nhan = (
                "lệnh"
                if k.ngon_ngu in {"bash", "powershell"}
                else "kết quả"
                if k.ngon_ngu == "text"
                else "mã"
            )
            hau_to = f" — {nhan} {n}" if len(code) > 1 else f" — {nhan}"
            slide_code(prs, muc.tieu_de + hau_to, k, muc, i)
        for i in [i for i, k in enumerate(muc.khoi) if k.loai == "table"][:1]:
            slide_bang(prs, muc.tieu_de + " — bảng đối chiếu", muc.khoi[i], muc, i)

    for n, tn in enumerate((noi_dung.trac_nghiem if noi_dung else []), start=1):
        slide_cau_hoi(prs, tn, n)
        slide_dap_an(prs, tn, n)

    muc_bt = next((m for m in bai.muc if m.tieu_de == MUC_BAI_TAP), None)
    if muc_bt:
        slide_bai_tap(prs, muc_bt)
    if noi_dung:
        slide_thao_luan(prs, noi_dung)
    slide_ket(prs, bai)
    slide_buoc_tiep(prs, bai)
    return prs


def main() -> int:
    parser = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument("--bai", nargs="*", help="số bài cần xuất, ví dụ: --bai 1 3 7")
    parser.add_argument("--dich", type=Path, default=DICH, help="thư mục xuất")
    tham_so = parser.parse_args()

    tep = sorted(NGUON.glob("session-*.md"))
    if tham_so.bai:
        chon = {f"{int(b):02d}" for b in tham_so.bai}
        tep = [t for t in tep if t.stem.split("-")[1] in chon]
    if not tep:
        print("Không tìm thấy bài nào trong docs/sessions/")
        return 1

    ngan_hang = doc_ngan_hang()
    tham_so.dich.mkdir(parents=True, exist_ok=True)
    for t in tep:
        bai = doc_bai(t)
        prs = dung_deck(bai, ngan_hang.get(bai.ma))
        ra = tham_so.dich / f"bai-{bai.ma}-{_ten_file(bai.tieu_de)}.pptx"
        prs.save(ra)
        duong = str(ra.relative_to(ROOT)) if ra.is_relative_to(ROOT) else str(ra)
        print(f"{duong}  ({len(prs.slides)} slide)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
